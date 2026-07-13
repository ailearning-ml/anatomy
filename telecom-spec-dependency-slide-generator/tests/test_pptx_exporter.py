"""Tests for PptxExporter and save_presentation."""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptx import Presentation as _Presentation

from src.main import build_analysis, parse_input_documents, save_presentation
from src.models import AnalysisOutput, GraphSummary, Slide
from src.pptx_exporter import PptxExporter

SAMPLE_PAYLOAD = [
    {
        "id": "ts-001",
        "title": "Core Exposure Spec",
        "domain": "core-network",
        "version": "1.0",
        "organization": "3GPP",
        "summary": "Defines service exposure.",
        "text": (
            "This specification depends on Subscriber Data Interface "
            "and includes workflow sequencing."
        ),
    },
    {
        "id": "ts-002",
        "title": "Subscriber Data Interface",
        "domain": "data-management",
        "version": "1.1",
        "organization": "ETSI",
        "summary": "Defines profile access.",
        "text": "This document defines interface expectations for subscriber data.",
    },
]


@pytest.fixture()
def sample_output() -> AnalysisOutput:
    documents = parse_input_documents(SAMPLE_PAYLOAD)
    return build_analysis(documents)


# ---------------------------------------------------------------------------
# PptxExporter.export
# ---------------------------------------------------------------------------


def test_export_creates_file(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "out.pptx"
    result = PptxExporter().export(sample_output, dest)
    assert result == dest
    assert dest.exists()


def test_export_file_is_non_empty(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "out.pptx"
    PptxExporter().export(sample_output, dest)
    assert dest.stat().st_size > 0


def test_export_creates_parent_directories(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "nested" / "deep" / "out.pptx"
    PptxExporter().export(sample_output, dest)
    assert dest.exists()


def test_export_slide_count(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "out.pptx"
    PptxExporter().export(sample_output, dest)
    prs = _Presentation(str(dest))
    # 1 cover slide + one slide per Slide object
    expected = 1 + len(sample_output.slides)
    assert len(prs.slides) == expected


def test_export_cover_title(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "out.pptx"
    custom_title = "My Custom Presentation"
    PptxExporter().export(sample_output, dest, title=custom_title)
    prs = _Presentation(str(dest))
    cover = prs.slides[0]
    title_shape = cover.shapes.title
    assert title_shape is not None
    assert title_shape.text == custom_title


def test_export_cover_custom_subtitle(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "out.pptx"
    PptxExporter().export(sample_output, dest, subtitle="My Subtitle")
    prs = _Presentation(str(dest))
    cover = prs.slides[0]
    placeholders = {ph.placeholder_format.idx: ph for ph in cover.placeholders}
    assert placeholders[1].text == "My Subtitle"


def test_export_content_slide_titles(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "out.pptx"
    PptxExporter().export(sample_output, dest)
    prs = _Presentation(str(dest))
    # Slides 1..N correspond to sample_output.slides
    for i, slide_data in enumerate(sample_output.slides, start=1):
        pptx_slide = prs.slides[i]
        assert pptx_slide.shapes.title is not None
        assert pptx_slide.shapes.title.text == slide_data.title


def test_export_minimal_output_no_slides(tmp_path: Path) -> None:
    """Exporter should still work when there are no content slides."""
    output = AnalysisOutput(
        slides=[],
        graph_summary=GraphSummary(
            total_documents=0,
            total_dependencies=0,
            critical_dependencies=0,
            circular_dependencies=0,
            interdomain_dependencies=0,
        ),
    )
    dest = tmp_path / "empty.pptx"
    PptxExporter().export(output, dest)
    prs = _Presentation(str(dest))
    assert len(prs.slides) == 1  # only the cover


def test_export_bullet_truncation(tmp_path: Path) -> None:
    """Content longer than _MAX_BULLETS should be truncated gracefully."""
    from src.pptx_exporter import _MAX_BULLETS, _TRUNCATION_NOTE

    many_bullets = [f"Bullet {i}" for i in range(_MAX_BULLETS + 5)]
    output = AnalysisOutput(
        slides=[Slide(title="Big Slide", purpose="Testing truncation", content=many_bullets)],
    )
    dest = tmp_path / "truncated.pptx"
    PptxExporter().export(output, dest)
    prs = _Presentation(str(dest))
    content_slide = prs.slides[1]
    all_text = "\n".join(
        ph.text_frame.text for ph in content_slide.placeholders if ph.has_text_frame
    )
    assert _TRUNCATION_NOTE in all_text


# ---------------------------------------------------------------------------
# save_presentation helper
# ---------------------------------------------------------------------------


def test_save_presentation_creates_file(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "via_helper.pptx"
    save_presentation(sample_output, dest)
    assert dest.exists()
    assert dest.stat().st_size > 0


def test_save_presentation_slide_count(tmp_path: Path, sample_output: AnalysisOutput) -> None:
    dest = tmp_path / "via_helper.pptx"
    save_presentation(sample_output, dest)
    prs = _Presentation(str(dest))
    assert len(prs.slides) == 1 + len(sample_output.slides)


# ---------------------------------------------------------------------------
# CLI --pptx-output integration
# ---------------------------------------------------------------------------


def test_cli_example_pptx_output(tmp_path: Path) -> None:
    from src.cli import main as cli_main

    json_out = tmp_path / "analysis.json"
    pptx_out = tmp_path / "deck.pptx"
    sys.argv = [
        "cli",
        "--example",
        "--output", str(json_out),
        "--pptx-output", str(pptx_out),
    ]
    result = cli_main()
    assert result == 0
    assert json_out.exists()
    assert pptx_out.exists()
    assert pptx_out.stat().st_size > 0


def test_cli_input_pptx_output(tmp_path: Path) -> None:
    import json
    from src.cli import main as cli_main

    input_file = tmp_path / "docs.json"
    input_file.write_text(json.dumps(SAMPLE_PAYLOAD), encoding="utf-8")
    json_out = tmp_path / "analysis.json"
    pptx_out = tmp_path / "deck.pptx"
    sys.argv = [
        "cli",
        "--input", str(input_file),
        "--output", str(json_out),
        "--pptx-output", str(pptx_out),
    ]
    result = cli_main()
    assert result == 0
    assert json_out.exists()
    assert pptx_out.exists()
    assert pptx_out.stat().st_size > 0


def test_cli_without_pptx_output_no_file_created(tmp_path: Path) -> None:
    """Omitting --pptx-output must not create any .pptx file."""
    from src.cli import main as cli_main

    json_out = tmp_path / "analysis.json"
    sys.argv = ["cli", "--example", "--output", str(json_out)]
    cli_main()
    assert not any(tmp_path.glob("*.pptx"))
