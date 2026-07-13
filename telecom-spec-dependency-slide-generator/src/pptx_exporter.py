"""PowerPoint export for telecom dependency analysis output."""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import AnalysisOutput, Slide

_MAX_BULLETS = 8
_TRUNCATION_NOTE = "Additional items omitted for brevity."


class PptxExporter:
    """Export an :class:`~src.models.AnalysisOutput` to a ``.pptx`` file."""

    def export(
        self,
        analysis: AnalysisOutput,
        destination: str | Path,
        *,
        title: str = "Telecom Dependency Anatomy",
        subtitle: Optional[str] = None,
    ) -> Path:
        """Create a PowerPoint presentation and save it to *destination*.

        Args:
            analysis: The structured analysis output to render.
            destination: File path for the generated ``.pptx``.
            title: Presentation title used on the cover slide.
            subtitle: Optional subtitle for the cover slide.  When *None*, a
                brief stats line derived from *analysis* is used instead.

        Returns:
            The resolved :class:`~pathlib.Path` of the saved file.
        """
        prs = Presentation()
        self._add_cover_slide(prs, title, subtitle, analysis)
        for slide in analysis.slides:
            self._add_content_slide(prs, slide)

        dest = Path(destination)
        dest.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(dest))
        return dest

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _add_cover_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: Optional[str],
        analysis: AnalysisOutput,
    ) -> None:
        layout = prs.slide_layouts[0]  # "Title Slide" layout
        slide = prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = title

        body_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_placeholder is not None:
            if subtitle:
                body_placeholder.text = subtitle
            else:
                summary = analysis.graph_summary
                if summary is not None:
                    body_placeholder.text = (
                        f"Documents: {summary.total_documents}  |  "
                        f"Dependencies: {summary.total_dependencies}  |  "
                        f"High-criticality: {summary.critical_dependencies}"
                    )

    def _add_content_slide(self, prs: Presentation, slide_data: Slide) -> None:
        layout = prs.slide_layouts[1]  # "Title and Content" layout
        slide = prs.slides.add_slide(layout)

        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = slide_data.title

        body_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_placeholder is None:
            return

        tf = body_placeholder.text_frame
        tf.clear()

        # Purpose line as the first paragraph (slightly smaller / italic feel)
        if slide_data.purpose:
            p = tf.paragraphs[0]
            p.text = slide_data.purpose
            if p.runs:
                p.runs[0].font.size = Pt(14)
                p.runs[0].font.italic = True

        bullets = list(slide_data.content)
        truncated = len(bullets) > _MAX_BULLETS
        if truncated:
            bullets = bullets[:_MAX_BULLETS]

        for item in bullets:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1

        if truncated:
            p = tf.add_paragraph()
            p.text = _TRUNCATION_NOTE
            p.level = 1
            if p.runs:
                p.runs[0].font.italic = True
